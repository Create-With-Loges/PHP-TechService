from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_ppt():
    prs = Presentation()

    # Define a theme or background function
    def apply_design(slide):
        # A simple "unique" design: a colored side bar
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 248, 255)  # Alice Blue background

        # Add a side shape
        left = Inches(0)
        top = Inches(0)
        width = Inches(1)
        height = Inches(7.5)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(44, 62, 80) # Dark Blue/Grey
        shape.line.fill.background()

    def add_title_slide(title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0] # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        apply_design(slide)
        
        # Adjust title/subtitle positions since we have a sidebar
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        # Position
        title.left = Inches(1.5)
        title.top = Inches(2)
        title.width = Inches(8)
        title.height = Inches(1.5)
        
        subtitle.left = Inches(1.5)
        subtitle.top = Inches(4)
        subtitle.width = Inches(8)
        subtitle.height = Inches(1)

        # Text and Font
        title.text = title_text
        for paragraph in title.text_frame.paragraphs:
            paragraph.font.size = Pt(40)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(44, 62, 80)

        subtitle.text = subtitle_text
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.font.size = Pt(24)
            paragraph.font.color.rgb = RGBColor(52, 73, 94)

    def add_content_slide(title_text, content_text_list):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        apply_design(slide)

        title = slide.shapes.title
        title.text = title_text
        title.left = Inches(1.5)
        title.top = Inches(0.5)
        title.width = Inches(8)
        
        # Set Title Font
        for paragraph in title.text_frame.paragraphs:
            paragraph.font.size = Pt(32)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(44, 62, 80)
            paragraph.alignment = PP_ALIGN.LEFT

        body_shape = slide.placeholders[1]
        body_shape.left = Inches(1.5)
        body_shape.top = Inches(1.5)
        body_shape.width = Inches(8)
        body_shape.height = Inches(5.5)
        
        tf = body_shape.text_frame
        tf.clear() # Clear default empty paragraph

        for item in content_text_list:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            
            # Check if it's a sub-item (indented in our list logic)
            if item.strip().startswith("-"):
                 p.level = 1
                 p.text = item.strip().lstrip("-").strip()

            p.space_after = Pt(10)
            p.font.size = Pt(20) # Safe size for content
            p.font.color.rgb = RGBColor(0, 0, 0)

    # 1. Title
    add_title_slide("Tech Service Support", "A Web Application for Hardware & Software Solutions")

    # 2. About the project
    about_points = [
        "Tech Service Support is a web application designed to help people fix their computer problems.",
        "It connects users with professional technicians for hardware and software issues.",
        "Users can register, find a service provider, and request help easily.",
        "The system helps manage requests and tracks the status of repairs.",
        "Key features include user dashboards, provider listings, and service tracking."
    ]
    add_content_slide("1. About the Project", about_points)

    # 3. Hardware & software specifications
    specs_points = [
        "Hardware Specifications:",
        " - RAM: 2 GB",
        " - ROM (Storage): 256 GB",
        " - Processor linked with Windows 7 32-bit",
        "",
        "Software Specifications:",
        " - Operating System: Windows 7 32-bit",
        " - Backend Language: PHP version 7.3.2",
        " - Database: MySQL version 5.10",
        " - Server Environment: XAMPP version 3.2.2"
    ]
    add_content_slide("2. Hardware & Software Specifications", specs_points)

    # 4. Problem definition
    problem_points = [
        "Difficulty in finding reliable technicians for specific computer problems.",
        "Lack of transparency in repair costs and service status.",
        "Time-consuming process of visiting multiple shops to find the right expert.",
        "No easy way to keep track of previous repairs and service history.",
        "Communication gaps between customers and service providers."
    ]
    add_content_slide("3. Problem Definition", problem_points)

    # 5. System study
    study_points = [
        "Existing System:",
        " - Manual Search: People have to physically go to shops.",
        " - No Records: Paper-based or verbal tracking of complaints.",
        " - Slow Process: Takes a long time to get updates.",
        "",
        "Proposed System:",
        " - Online Platform: specific website to login and request service.",
        " - Digital Records: All data is saved in a database (MySQL).",
        " - Real-time Updates: specialized dashboard to view status.",
        " - Categories: Clear separation of Hardware and Software services."
    ]
    add_content_slide("4. System Study", study_points)

    # 6. Proposed system (More details)
    proposed_points = [
        "User Module: Users can sign up, log in, view providers, and send requests.",
        "Provider Module: Technicians can manage their profiles and accept/reject requests.",
        "Admin Module: Oversees the entire system operations.",
        "Status Tracking: Requests move from 'Pending' to 'In Progress' to 'Completed'.",
        "Security: Secure login for all users to protect personal data and service history."
    ]
    add_content_slide("5. Proposed System", proposed_points)

    # 7. Dataflow diagram (Visual Representation)
    slide_layout = prs.slide_layouts[5] # Title Only
    slide = prs.slides.add_slide(slide_layout)
    apply_design(slide)
    
    title = slide.shapes.title
    title.text = "6. Dataflow Diagram"
    title.left = Inches(1.5)
    title.top = Inches(0.5)
    
    # Set Title Font
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.size = Pt(32)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(44, 62, 80)
        paragraph.alignment = PP_ALIGN.LEFT

    # Draw a simple flow
    # User -> Request -> System -> Provider
    
    # 1. User
    left = Inches(2)
    top = Inches(2)
    width = Inches(1.5)
    height = Inches(1)
    shape_user = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape_user.text = "User\n(Registers/Logins)"
    
    # Arrow 1
    left = Inches(3.5)
    top = Inches(2.5)
    width = Inches(1)
    height = Inches(0.2)
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)

    # 2. System/Request
    left = Inches(4.5)
    top = Inches(2)
    width = Inches(2)
    height = Inches(1)
    shape_system = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, left, top, width, height)
    shape_system.text = "System\n(Process Request)"

    # Arrow 2
    left = Inches(6.5)
    top = Inches(2.5)
    width = Inches(1)
    height = Inches(0.2)
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)

    # 3. Provider
    left = Inches(7.5)
    top = Inches(2)
    width = Inches(1.5)
    height = Inches(1)
    shape_provider = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape_provider.text = "Service Provider\n(Accepts/Fixes)"

    # Status Update loop (Provider -> System -> User)
    # Arrow 3 (Back)
    left = Inches(6.5)
    top = Inches(3.5)
    width = Inches(1)
    height = Inches(0.2)
    arrow3 = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, left, top, width, height)
    
    # Status Text
    left = Inches(4.5)
    top = Inches(3.5)
    width = Inches(2)
    height = Inches(0.5)
    status_box = slide.shapes.add_textbox(left, top, width, height)
    status_box.text_frame.text = "Update Status"
    status_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Arrow 4 (Back to User)
    left = Inches(3.5)
    top = Inches(3.5)
    width = Inches(1)
    height = Inches(0.2)
    arrow4 = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, left, top, width, height)

    prs.save('Tech_Service_Support_v2.pptx')
    print("PPT generated successfully with updated fonts as v2.")

if __name__ == "__main__":
    create_ppt()
