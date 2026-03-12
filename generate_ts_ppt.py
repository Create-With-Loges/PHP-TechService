from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def set_cell_border(cell):
    """Adds a solid black border to a table cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for border_name in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        ln = tcPr.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{border_name.split(':')[1]}")
        if ln is None:
            # Add line element
            from pptx.oxml.xmlchemy import OxmlElement
            ln = OxmlElement(border_name)
            ln.set('w', '12700') # 1 pt width
            ln.set('cmpd', 'sng')
            solidFill = OxmlElement('a:solidFill')
            srgbClr = OxmlElement('a:srgbClr')
            srgbClr.set('val', '000000') # Black color
            solidFill.append(srgbClr)
            ln.append(solidFill)
            tcPr.append(ln)

def style_table(table, data, headers):
    """Styles a table to be black and white with borders."""
    # Set headers
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        set_cell_border(cell)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(230, 230, 230) # Light gray header background
        
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(14)
                run.font.color.rgb = RGBColor(0, 0, 0) # Black text

    # Set data
    for row_idx, row_data in enumerate(data):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = text
            set_cell_border(cell)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) # White background
            
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = RGBColor(0, 0, 0) # Black text

# Create presentation
prs = Presentation()

# --- Slide 1: Database File Table Specification ---
slide_layout = prs.slide_layouts[5] # blank slide with title
slide1 = prs.slides.add_slide(slide_layout)
title_shape = slide1.shapes.title
title_shape.text = "Database File Table Specification"

# Common Headers
headers = ['Column Name', 'Data Type', 'Constraints', 'Description']

# Admins Table Title
txBox = slide1.shapes.add_textbox(Pt(50), Pt(100), Pt(200), Pt(30))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Admins Table"
p.font.bold = True
p.font.size = Pt(16)

# Admins Purpose
txBox_p1 = slide1.shapes.add_textbox(Pt(50), Pt(130), Pt(600), Pt(20))
tf_p1 = txBox_p1.text_frame
p_p1 = tf_p1.paragraphs[0]
p_p1.text = "Purpose: Stores overarching administrator credentials with full system access."
p_p1.font.size = Pt(12)
p_p1.font.italic = True

# Add Admins Table
admins_data = [
    ['id', 'INT', 'PRIMARY KEY, AUTO_INCREMENT', 'Unique identifier for admin records'],
    ['username', 'VARCHAR(50)', 'NOT NULL', 'Administrator login alias'],
    ['password', 'VARCHAR(255)', 'NOT NULL', 'Encrypted/Secure login password'],
    ['email', 'VARCHAR(100)', 'None', 'Administrator contact email'],
    ['created_at', 'TIMESTAMP', 'DEFAULT CURRENT_TIMESTAMP', 'Timestamp of record generation']
]
rows1, cols1 = len(admins_data) + 1, 4
left1, top1, width1, height1 = Pt(50), Pt(160), Pt(600), Pt(100)
table1 = slide1.shapes.add_table(rows1, cols1, left1, top1, width1, height1).table
table1.columns[0].width = Pt(120)
table1.columns[1].width = Pt(140)
table1.columns[2].width = Pt(180)
table1.columns[3].width = Pt(220)

style_table(table1, admins_data, headers)

# Users Table Title
txBox2 = slide1.shapes.add_textbox(Pt(50), Pt(320), Pt(200), Pt(30))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "Users Table"
p2.font.bold = True
p2.font.size = Pt(16)

# Users Purpose
txBox_p2 = slide1.shapes.add_textbox(Pt(50), Pt(350), Pt(600), Pt(20))
tf_p2 = txBox_p2.text_frame
p_p2 = tf_p2.paragraphs[0]
p_p2.text = "Purpose: Holds the registered details of clients seeking technical assistance."
p_p2.font.size = Pt(12)
p_p2.font.italic = True

# Add Users Table
users_data = [
    ['id', 'INT', 'PRIMARY KEY, AUTO_INCREMENT', 'Client unique identifier'],
    ['username', 'VARCHAR(50)', 'NOT NULL', 'Client display/login name'],
    ['password', 'VARCHAR(255)', 'NOT NULL', 'Client secure password'],
    ['email', 'VARCHAR(100)', 'None', 'Client registered email address'],
    ['contact', 'VARCHAR(20)', 'None', 'Client mobile phone number'],
    ['model/company', 'VARCHAR(100)', 'None', 'Device model and brand details']
]
rows2, cols2 = len(users_data) + 1, 4
left2, top2, width2, height2 = Pt(50), Pt(380), Pt(600), Pt(120)
table2 = slide1.shapes.add_table(rows2, cols2, left2, top2, width2, height2).table
for c in range(4): table2.columns[c].width = table1.columns[c].width

style_table(table2, users_data, headers)


# --- Slide 2: Database File Table Specification (Continued) ---
slide2 = prs.slides.add_slide(slide_layout)
title_shape2 = slide2.shapes.title
title_shape2.text = "Database File Table Specification (Cont.)"

# Providers Table Title
txBox3 = slide2.shapes.add_textbox(Pt(50), Pt(80), Pt(200), Pt(30))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "Providers Table"
p3.font.bold = True
p3.font.size = Pt(16)

# Providers Purpose
txBox_p3 = slide2.shapes.add_textbox(Pt(50), Pt(110), Pt(600), Pt(20))
tf_p3 = txBox_p3.text_frame
p_p3 = tf_p3.paragraphs[0]
p_p3.text = "Purpose: Registers service technicians along with their business identifiers and skills."
p_p3.font.size = Pt(12)
p_p3.font.italic = True

# Add Providers Table
providers_data = [
    ['id', 'INT', 'PRIMARY KEY, AUTO_INCREMENT', 'Technician unique ID'],
    ['username/password', 'VARCHAR', 'NOT NULL', 'Technician login credentials'],
    ['email/contact', 'VARCHAR', 'None', 'Technician contact pathways'],
    ['shop_reg_num', 'VARCHAR(100)', 'None', 'Verified business registration number'],
    ['service_type', 'VARCHAR(50)', 'None', 'Core category (e.g., Hardware, Software)'],
    ['short_info', 'TEXT', 'None', 'Brief professional bio or skill summary']
]
rows3, cols3 = len(providers_data) + 1, 4
left3, top3, width3, height3 = Pt(50), Pt(140), Pt(600), Pt(120)
table3 = slide2.shapes.add_table(rows3, cols3, left3, top3, width3, height3).table
for c in range(4): table3.columns[c].width = table1.columns[c].width

style_table(table3, providers_data, headers)

# Requests Table Title
txBox4 = slide2.shapes.add_textbox(Pt(50), Pt(320), Pt(200), Pt(30))
tf4 = txBox4.text_frame
p4 = tf4.paragraphs[0]
p4.text = "Requests Table"
p4.font.bold = True
p4.font.size = Pt(16)

# Requests Purpose
txBox_p4 = slide2.shapes.add_textbox(Pt(50), Pt(350), Pt(600), Pt(20))
tf_p4 = txBox_p4.text_frame
p_p4 = tf_p4.paragraphs[0]
p_p4.text = "Purpose: Manages workflow tickets linking Users to Providers."
p_p4.font.size = Pt(12)
p_p4.font.italic = True

# Add Requests Table
requests_data = [
    ['id', 'INT', 'PRIMARY KEY, AUTO_INCREMENT', 'Unique service ticket ID'],
    ['user_id', 'INT', 'FOREIGN KEY (users)', 'Referential link to the requesting client'],
    ['provider_id', 'INT', 'FOREIGN KEY (providers)', 'Referential link to assigned technician'],
    ['details', 'TEXT', 'None', 'Specific breakdown of the issue'],
    ['status', 'ENUM', "DEFAULT 'Pending'", 'Progress state (Pending, Completed, etc.)']
]
rows4, cols4 = len(requests_data) + 1, 4
left4, top4, width4, height4 = Pt(50), Pt(380), Pt(600), Pt(100)
table4 = slide2.shapes.add_table(rows4, cols4, left4, top4, width4, height4).table
for c in range(4): table4.columns[c].width = table1.columns[c].width

style_table(table4, requests_data, headers)


# --- Slide 3: Module Specification ---
slide_layout_txt = prs.slide_layouts[1] # Title and Content
slide3 = prs.slides.add_slide(slide_layout_txt)
title_shape3 = slide3.shapes.title
title_shape3.text = "Module Specification"

body_shape = slide3.shapes.placeholders[1]
tf_body = body_shape.text_frame
tf_body.clear()

def add_detailed_module(tf, num_name, subtitle, desc):
    # Heading (font size 18)
    p1 = tf.add_paragraph()
    p1.text = num_name
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.space_before = Pt(5)
    
    # Sub heading (font size 16)
    p2 = tf.add_paragraph()
    p2.text = f"[{subtitle}]"
    p2.font.size = Pt(16)
    p2.font.italic = True
    p2.font.color.rgb = RGBColor(100, 100, 100)
    
    # Body content (font size 14)
    p3 = tf.add_paragraph()
    p3.text = desc
    p3.font.size = Pt(14)
    p3.space_after = Pt(10)

add_detailed_module(tf_body, "1. User Interface (Client Portal)", "Service Request Module", 
    "Acting as the primary entry point for individuals seeking technical solutions, this module allows registered users to outline their specific issues. They can input device demographics like model and company, review available technicians, initiate a direct service request, and actively monitor the live 'status' of their pending queries.")

add_detailed_module(tf_body, "2. Provider Dashboard", "Technician Workflow Module", 
    "Serves as the dedicated workspace for registered technicians (Providers). It allows them to showcase their professional competencies, service types, and external portfolios. Critically, it acts as an incoming queue where they can selectively review, accept, or reject incoming hardware/software service tickets dispatched by users.")

add_detailed_module(tf_body, "3. Administrative Control Panel", "System Oversight Module", 
    "Designed for overarching platform governors, this segment permits absolute visibility across all transactions. Administrators have the authoritative capacity to manage foundational data layers, oversee user-to-provider interactions, review system-wide contact queries, and enforce cascading database referential integrity checks.")

# --- Slide 4: Entity-Relationship Diagram ---
slide_layout_blank = prs.slide_layouts[5] # blank slide with title
slide4 = prs.slides.add_slide(slide_layout_blank)
title_shape4 = slide4.shapes.title
title_shape4.text = "Entity-Relationship (ER) Diagram"

# Add ER Diagram Image
img_path = r"C:\Users\ASUS\.gemini\antigravity\brain\f4c97a84-f519-4fcc-841c-35c2d0e2b9ca\ts_er_diagram_1772128673377.png"
left_img = Pt(50)
top_img = Pt(120)
height_img = Pt(400) # scale height, width auto-adjusts
slide4.shapes.add_picture(img_path, left_img, top_img, height=height_img)

# Save the presentation
prs.save("TS_Project_Analysis.pptx")
print("Presentation generated successfully!")
